from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
from typing import Dict, Any
import os
import tempfile
from datetime import datetime

from .database import SessionLocal
from .models import Project, Content, Refinement, User
from .schemas import (
    GenerateContentRequest,
    RefineContentRequest,
    FeedbackRequest,
    AIGenerateOutlineRequest
)
from .auth.routes import get_current_user
from .services.gemini_service import generate_content, refine_content, generate_outline
from docx import Document
from pptx import Presentation
from pptx.util import Pt

router = APIRouter(prefix="/projects", tags=["Documents"])


def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


@router.post("/{project_id}/generate")
async def generate_project_content(
    project_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Generate content for all sections/slides in a project"""
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == current_user.id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    if not project.outline:
        raise HTTPException(status_code=400, detail="Project outline is empty")
    
    if not project.topic:
        raise HTTPException(status_code=400, detail="Project topic is required")
    
    generated_content = {}
    
    for item in project.outline:
        section_id = item.get("id") or item.get("title", "").lower().replace(" ", "_")
        section_title = item.get("title", item.get("name", ""))
        
        # Generate content using Gemini
        content_text = generate_content(
            topic=project.topic,
            section_title=section_title,
            doc_type=project.doc_type or "docx"
        )
        
        generated_content[section_id] = {
            "title": section_title,
            "content": content_text,
            "generated_at": datetime.now().isoformat()
        }
        
        # Store in Content table
        content_obj = Content(
            project_id=project.id,
            section_id=section_id,
            text=content_text
        )
        db.add(content_obj)
    
    # Update project content
    project.content = generated_content
    db.commit()
    db.refresh(project)
    
    return {
        "message": "Content generated successfully",
        "content": generated_content
    }


@router.post("/{project_id}/generate_section")
async def generate_single_section(
    project_id: int,
    request: GenerateContentRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Generate content for a single section/slide"""
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == current_user.id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Find section in outline
    section = None
    for item in project.outline:
        if item.get("id") == request.section_id or item.get("title", "").lower().replace(" ", "_") == request.section_id:
            section = item
            break
    
    if not section:
        raise HTTPException(status_code=404, detail="Section not found")
    
    section_title = section.get("title", section.get("name", ""))
    
    # Get existing content if any
    existing_content = None
    if project.content and request.section_id in project.content:
        existing_content = project.content[request.section_id].get("content", "")
    
    # Generate content
    content_text = generate_content(
        topic=project.topic or "",
        section_title=section_title,
        doc_type=project.doc_type or "docx",
        existing_content=existing_content
    )
    
    # Update project content
    if not project.content:
        project.content = {}
    
    project.content[request.section_id] = {
        "title": section_title,
        "content": content_text,
        "generated_at": datetime.now().isoformat()
    }
    
    # Store in Content table
    content_obj = Content(
        project_id=project.id,
        section_id=request.section_id,
        text=content_text
    )
    db.add(content_obj)
    db.commit()
    
    return {
        "section_id": request.section_id,
        "content": content_text
    }


@router.post("/{project_id}/refine")
async def refine_section_content(
    project_id: int,
    request: RefineContentRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Refine content for a specific section/slide"""
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == current_user.id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    if not project.content or request.section_id not in project.content:
        raise HTTPException(status_code=404, detail="Section content not found")
    
    original_content = project.content[request.section_id].get("content", "")
    section_title = project.content[request.section_id].get("title", "")
    
    # Refine using Gemini
    refined_text = refine_content(
        original_content=original_content,
        refinement_prompt=request.refinement_prompt,
        topic=project.topic or "",
        section_title=section_title
    )
    
    # Update content
    project.content[request.section_id]["content"] = refined_text
    project.content[request.section_id]["refined_at"] = datetime.now().isoformat()
    
    # Store refinement history
    if not project.refinement_history:
        project.refinement_history = []
    
    project.refinement_history.append({
        "section_id": request.section_id,
        "prompt": request.refinement_prompt,
        "timestamp": datetime.now().isoformat()
    })
    
    # Store in Refinement table
    refinement_obj = Refinement(
        project_id=project.id,
        section_id=request.section_id,
        prompt=request.refinement_prompt,
        updated_text=refined_text
    )
    db.add(refinement_obj)
    
    # Update Content table
    content_obj = db.query(Content).filter(
        Content.project_id == project.id,
        Content.section_id == request.section_id
    ).order_by(Content.updated_at.desc()).first()
    
    if content_obj:
        content_obj.text = refined_text
    else:
        content_obj = Content(
            project_id=project.id,
            section_id=request.section_id,
            text=refined_text
        )
        db.add(content_obj)
    
    db.commit()
    db.refresh(project)
    
    return {
        "section_id": request.section_id,
        "refined_content": refined_text
    }


@router.post("/{project_id}/feedback")
async def submit_feedback(
    project_id: int,
    request: FeedbackRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Submit feedback (like/dislike/comment) for a section"""
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == current_user.id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    if not project.feedback:
        project.feedback = {}
    
    if request.section_id not in project.feedback:
        project.feedback[request.section_id] = {}
    
    if request.like is not None:
        project.feedback[request.section_id]["like"] = request.like
    
    if request.comment:
        if "comments" not in project.feedback[request.section_id]:
            project.feedback[request.section_id]["comments"] = []
        project.feedback[request.section_id]["comments"].append({
            "comment": request.comment,
            "timestamp": datetime.now().isoformat()
        })
    
    db.commit()
    db.refresh(project)
    
    return {"message": "Feedback submitted successfully"}


@router.post("/{project_id}/ai-outline")
async def generate_ai_outline(
    project_id: int,
    request: AIGenerateOutlineRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Generate outline using AI"""
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == current_user.id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Generate outline using Gemini
    outline_titles = generate_outline(request.topic, request.doc_type)
    
    # Format outline
    outline = []
    for idx, title in enumerate(outline_titles):
        section_id = f"section_{idx + 1}" if request.doc_type == "docx" else f"slide_{idx + 1}"
        outline.append({
            "id": section_id,
            "title": title,
            "order": idx + 1
        })
    
    project.outline = outline
    project.topic = request.topic
    project.doc_type = request.doc_type
    
    db.commit()
    db.refresh(project)
    
    return {
        "outline": outline,
        "message": "Outline generated successfully"
    }


@router.get("/{project_id}/export/docx")
async def export_docx(
    project_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Export project as Word document"""
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == current_user.id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    if project.doc_type != "docx":
        raise HTTPException(status_code=400, detail="Project is not a Word document")
    
    if not project.content:
        raise HTTPException(status_code=400, detail="No content to export")
    
    # Create Word document
    doc = Document()
    
    # Add title
    doc.add_heading(project.title, 0)
    
    # Add topic if available
    if project.topic:
        doc.add_paragraph(f"Topic: {project.topic}")
        doc.add_paragraph("")  # Empty line
    
    # Add sections
    for item in project.outline:
        section_id = item.get("id") or item.get("title", "").lower().replace(" ", "_")
        section_title = item.get("title", item.get("name", ""))
        
        # Add heading
        doc.add_heading(section_title, level=1)
        
        # Add content
        if section_id in project.content:
            content_text = project.content[section_id].get("content", "")
            doc.add_paragraph(content_text)
        
        doc.add_paragraph("")  # Empty line between sections
    
    # Save to temporary file
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
    doc.save(temp_file.name)
    temp_file.close()
    
    return FileResponse(
        temp_file.name,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        filename=f"{project.title.replace(' ', '_')}.docx"
    )


@router.get("/{project_id}/export/pptx")
async def export_pptx(
    project_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Export project as PowerPoint presentation"""
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == current_user.id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    if project.doc_type != "pptx":
        raise HTTPException(status_code=400, detail="Project is not a PowerPoint presentation")
    
    if not project.content:
        raise HTTPException(status_code=400, detail="No content to export")
    
    # Create PowerPoint presentation
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = project.title
    if project.topic:
        subtitle.text = project.topic
    
    # Add content slides
    for item in project.outline:
        section_id = item.get("id") or item.get("title", "").lower().replace(" ", "_")
        section_title = item.get("title", item.get("name", ""))
        
        # Use Title and Content layout
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = section_title
        
        tf = body_shape.text_frame
        tf.text = ""  # Clear default text
        
        if section_id in project.content:
            content_text = project.content[section_id].get("content", "")
            # Split into bullet points
            paragraphs = content_text.split('\n')
            first = True
            for para in paragraphs:
                if para.strip():
                    if first:
                        tf.text = para.strip()
                        first = False
                    else:
                        p = tf.add_paragraph()
                        p.text = para.strip()
                        p.level = 0
                        p.font.size = Pt(14)
    
    # Save to temporary file
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(temp_file.name)
    temp_file.close()
    
    return FileResponse(
        temp_file.name,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{project.title.replace(' ', '_')}.pptx"
    )

